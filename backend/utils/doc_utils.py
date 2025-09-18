import fitz # PyMuPDF
from pptx import Presentation   

def extract_text_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    return " ".join([page.get_text() for page in doc])

def extract_text_ppt(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return " ".join(text_runs)