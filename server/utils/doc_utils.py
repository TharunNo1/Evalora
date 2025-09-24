# import fitz # PyMuPDF
from pptx import Presentation   
import os
import logging

# ---------- Third-party imports (install as needed) ----------
from docx import Document               
from pptx import Presentation           
from email import policy
from email.parser import BytesParser
from pydub import AudioSegment          
import speech_recognition as sr         
# from moviepy.editor import AudioFileClip 
import pdfplumber

def extract_text_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_text_ppt(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return " ".join(text_runs)

def extract_text_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs).strip()


def extract_text_doc(file_path: str) -> str:
    """Extract text from legacy DOC using textract or antiword."""
    import textract  # requires `textract` + system dependencies
    return textract.process(file_path).decode("utf-8").strip()


def extract_text_txt(file_path: str) -> str:
    """Extract text from a plain TXT file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_text_eml(file_path: str) -> str:
    """Extract text from an EML (email) file."""
    with open(file_path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    return msg.get_body(preferencelist=('plain', 'html')).get_content().strip()


# ---------- Audio/Video to text ----------

def _audio_to_text(audio_path: str) -> str:
    """Convert an audio file to text using SpeechRecognition."""
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio = recognizer.record(source)
    return recognizer.recognize_google(audio)


def extract_text_audio(file_path: str) -> str:
    """
    Handle mp3, wav, m4a audio. Converts to WAV (if needed)
    and then transcribes to text.
    """
    # Convert to wav if not already
    if not file_path.lower().endswith(".wav"):
        audio = AudioSegment.from_file(file_path)
        temp_wav = file_path + ".tmp.wav"
        audio.export(temp_wav, format="wav")
        text = _audio_to_text(temp_wav)
        os.remove(temp_wav)
    else:
        text = _audio_to_text(file_path)
    return text


# ---------- Dispatcher ----------

def extract_text(file_path: str) -> str:
    """
    Route the file to the correct extractor based on extension.
    Returns an empty string if no extractor matches.
    """
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return extract_text_pdf(file_path)
        if ext in {".ppt", ".pptx"}:
            return extract_text_ppt(file_path)
        if ext in {".docx"}:
            return extract_text_docx(file_path)
        if ext in {".doc"}:
            return extract_text_doc(file_path)
        if ext == ".txt":
            return extract_text_txt(file_path)
        if ext == ".eml":
            return extract_text_eml(file_path)
        if ext in {".mp3", ".wav", ".m4a"}:
            return extract_text_audio(file_path)
        # if ext == ".mp4":
        #     return extract_text_mp4(file_path)
    except Exception as exc:
        logger.exception("Failed to extract text from %s: %s", file_path, exc)

    logger.warning("Unsupported file type: %s", file_path)
    return ""
